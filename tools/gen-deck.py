#!/usr/bin/env python3
"""Generate the SynapseOS "AI-native OS" pitch deck.

Usage:  python tools/gen-deck.py [output.pptx]

Needs python-pptx. The palette matches the distro's Catppuccin Macchiato theme
so the deck reads as part of the same product as the ISO branding.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "archiso/airootfs/etc/calamares/branding/synapseos/logo.png"

# ---------------------------------------------------------------- palette ----
CRUST = RGBColor(0x18, 0x19, 0x26)
MANTLE = RGBColor(0x1E, 0x20, 0x30)
BASE = RGBColor(0x24, 0x27, 0x3A)
SURFACE0 = RGBColor(0x36, 0x3A, 0x4F)
SURFACE1 = RGBColor(0x49, 0x4D, 0x64)
OVERLAY = RGBColor(0x6E, 0x73, 0x8D)
SUBTEXT = RGBColor(0xA5, 0xAD, 0xCB)
TEXT = RGBColor(0xCA, 0xD3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MAUVE = RGBColor(0xC6, 0xA0, 0xF6)
BLUE = RGBColor(0x8A, 0xAD, 0xF4)
SAPPHIRE = RGBColor(0x7D, 0xC4, 0xE4)
TEAL = RGBColor(0x8B, 0xD5, 0xCA)
GREEN = RGBColor(0xA6, 0xDA, 0x95)
YELLOW = RGBColor(0xEE, 0xD4, 0x9F)
PEACH = RGBColor(0xF5, 0xA9, 0x7F)
RED = RGBColor(0xED, 0x87, 0x96)

FONT = "Inter"
MONO = "DejaVu Sans Mono"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.72)          # page margin
CW = W - 2 * M            # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# ---------------------------------------------------------------- helpers ----
def _noline(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    return _noline(s)


def card(slide, x, y, w, h, fill=MANTLE, accent=None, radius=0.055):
    """Rounded panel, optionally with a coloured spine on the left edge."""
    s = rect(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, radius)
    if accent is not None:
        rect(slide, x, y + Inches(0.16), Inches(0.045), h - Inches(0.32),
             accent, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5)
    return s


def tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, first=False):
    """lines: list of dicts -> text, size, color, bold, font, space_after, line, indent, bullet."""
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and first) else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        p.space_after = Pt(ln.get("space_after", 0))
        if ln.get("line"):
            p.line_spacing = ln["line"]
        txt = ln["t"]
        if ln.get("bullet"):
            txt = "•   " + txt
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ln.get("font", FONT)
        f.size = Pt(ln.get("size", 14))
        f.bold = ln.get("bold", False)
        f.color.rgb = ln.get("c", TEXT)
    return tf


def slide_base(title, kicker=None, footer=True):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, BASE)
    # top hairline accent
    rect(s, 0, 0, W, Inches(0.055), MAUVE)

    y = Inches(0.62)
    if kicker:
        write(tbox(s, M, y, CW, Inches(0.26)),
              [{"t": kicker.upper(), "size": 10.5, "c": MAUVE, "bold": True}], first=True)
        y += Inches(0.34)
    write(tbox(s, M, y, CW, Inches(0.62)),
          [{"t": title, "size": 31, "c": WHITE, "bold": True}], first=True)
    rect(s, M, y + Inches(0.72), Inches(1.5), Inches(0.035), SURFACE1)

    if footer:
        _page["n"] += 1
        rect(s, 0, H - Inches(0.42), W, Inches(0.42), MANTLE)
        write(tbox(s, M, H - Inches(0.31), CW / 2, Inches(0.2)),
              [{"t": "SynapseOS  ·  AI-native operating system",
                "size": 9, "c": OVERLAY}], first=True)
        write(tbox(s, M + CW / 2, H - Inches(0.31), CW / 2, Inches(0.2)),
              [{"t": f"{_page['n']:02d}", "size": 9, "c": OVERLAY,
                "align": PP_ALIGN.RIGHT}], first=True)
    return s, Inches(1.62)


def chip(slide, x, y, text, color, w=None):
    tw = w or Inches(0.11 * len(text) + 0.34)
    rect(slide, x, y, tw, Inches(0.3), SURFACE0, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5)
    write(tbox(slide, x, y + Inches(0.065), tw, Inches(0.2)),
          [{"t": text, "size": 10, "c": color, "bold": True,
            "align": PP_ALIGN.CENTER}], first=True)
    return x + tw + Inches(0.12)


def grid(slide, y, items, cols=3, h=Inches(1.72), gap=Inches(0.22),
         num=False, body_size=11.5):
    """items: (accent, heading, body)."""
    cw = int((CW - gap * (cols - 1)) / cols)
    for i, (accent, head, body) in enumerate(items):
        col, row = i % cols, i // cols
        x = M + col * (cw + gap)
        yy = y + row * (h + gap)
        card(slide, x, yy, cw, h, MANTLE, accent)
        tf = tbox(slide, x + Inches(0.3), yy + Inches(0.24),
                  cw - Inches(0.52), h - Inches(0.44))
        lines = []
        if num:
            lines.append({"t": f"{i + 1:02d}", "size": 10, "c": accent,
                          "bold": True, "space_after": 5, "font": MONO})
        lines.append({"t": head, "size": 14.5, "c": WHITE, "bold": True,
                      "space_after": 7, "line": 1.05})
        lines.append({"t": body, "size": body_size, "c": SUBTEXT, "line": 1.32})
        write(tf, lines, first=True)


def rows(slide, y, items, h=Inches(0.86), gap=Inches(0.13), label_w=Inches(3.5)):
    """items: (accent, label, detail)."""
    for i, (accent, label, detail) in enumerate(items):
        yy = y + i * (h + gap)
        card(slide, M, yy, CW, h, MANTLE, accent)
        write(tbox(slide, M + Inches(0.3), yy, label_w, h, MSO_ANCHOR.MIDDLE),
              [{"t": label, "size": 13.5, "c": WHITE, "bold": True, "line": 1.1}],
              first=True)
        rect(slide, M + Inches(0.3) + label_w, yy + Inches(0.18),
             Inches(0.012), h - Inches(0.36), SURFACE1)
        write(tbox(slide, M + Inches(0.55) + label_w, yy,
                   CW - label_w - Inches(0.9), h, MSO_ANCHOR.MIDDLE),
              [{"t": detail, "size": 11.5, "c": SUBTEXT, "line": 1.3}], first=True)


# ============================================================== 01  title ====
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, CRUST)
rect(s, 0, 0, Inches(0.09), H, MAUVE)
# faint geometry, right-anchored inside the margin
for i in range(3):
    w = Inches(3.9 - 0.55 * i)
    rect(s, W - M - w, Inches(1.1) + Inches(0.42 * i), w, Inches(0.02), SURFACE0)

if LOGO.exists():
    s.shapes.add_picture(str(LOGO), M, Inches(1.28), Inches(1.02), Inches(1.02))

write(tbox(s, M, Inches(2.62), Inches(11.0), Inches(0.34)),
      [{"t": "AN OPERATING SYSTEM WITH A MIND", "size": 12, "c": MAUVE,
        "bold": True}], first=True)
write(tbox(s, M, Inches(3.06), Inches(11.4), Inches(1.5)),
      [{"t": "SynapseOS", "size": 76, "c": WHITE, "bold": True}], first=True)
write(tbox(s, M, Inches(4.42), Inches(9.6), Inches(1.1)),
      [{"t": "An AI-native Linux distribution where the assistant is a system "
             "service — it sees your messages, apps and processes, and acts "
             "on them under permission you control.",
        "size": 16.5, "c": SUBTEXT, "line": 1.42}], first=True)
rect(s, M, Inches(5.72), Inches(2.0), Inches(0.035), MAUVE)

x = M
for label, cl in (("Arch + Plasma 6", BLUE), ("Local-first inference", TEAL),
                  ("Capability-gated actions", GREEN), ("Open connector SDK", YELLOW)):
    x = chip(s, x, Inches(6.12), label, cl)

write(tbox(s, M, Inches(6.86), CW, Inches(0.22)),
      [{"t": "Concept & architecture deck", "size": 10, "c": OVERLAY}], first=True)

# ============================================================ 02  problem ====
s, y = slide_base("AI is the smartest thing on your machine, and the most blind",
                  "The problem")
write(tbox(s, M, y, Inches(10.6), Inches(0.5)),
      [{"t": "Today's assistant lives in a tab. It cannot see what you see, and it "
             "cannot touch anything.", "size": 14.5, "c": SUBTEXT, "line": 1.35}],
      first=True)
grid(s, y + Inches(0.72), [
    (RED, "You are the integration layer",
     "Copy, paste, screenshot, re-explain. Every task starts by rebuilding context "
     "the machine already has."),
    (PEACH, "Context is shattered",
     "Telegram, WhatsApp, Slack, Gmail, Discord. Six inboxes, no shared timeline, "
     "nothing correlated."),
    (YELLOW, "Attention has no filter",
     "Hundreds of notifications a day, ranked by arrival time instead of by what "
     "actually needs you."),
    (BLUE, "It cannot act",
     "No assistant can close the runaway process draining your battery, even when "
     "it can name it."),
    (MAUVE, "The OS knows, and tells no one",
     "/proc, D-Bus and the window manager expose rich state. No reasoning layer "
     "ever consumes it."),
    (TEAL, "Bolt-on integrations rot",
     "Per-app plugins duplicate auth, permissions and memory. Nothing composes "
     "across them."),
], cols=3, h=Inches(1.86))

# ============================================================= 03  thesis ====
s, y = slide_base("Put the assistant inside the OS, not on top of it", "The idea")
card(s, M, y, CW, Inches(2.05), MANTLE, MAUVE)
write(tbox(s, M + Inches(0.52), y + Inches(0.3), CW - Inches(1.1), Inches(1.5)),
      [{"t": "SynapseOS treats intelligence as a system service. One privileged "
             "reasoning layer perceives the whole machine — messages, windows, "
             "processes, files — and can act on it through a permission model "
             "the user owns.", "size": 21, "c": WHITE, "line": 1.34}], first=True)

grid(s, y + Inches(2.4), [
    (BLUE, "One brain, not twelve plugins",
     "A single core with shared memory, identity and policy. Apps plug into it "
     "instead of each shipping their own assistant."),
    (TEAL, "Perception is first-class",
     "Message streams, app state and the process table are OS-level inputs, "
     "sampled continuously and normalised."),
    (GREEN, "Action is capability-gated",
     "Every effect — send, launch, kill, throttle — is a scoped, revocable, "
     "audited capability call. Denied by default."),
], cols=3, h=Inches(1.86))

# ========================================================= 04  experience ====
s, y = slide_base("What it feels like to use", "The experience")
rows(s, y, [
    (SAPPHIRE, "“What did I miss?”",
     "One briefing across every account: three messages need a reply, one is "
     "blocking someone, the rest are noise. Grouped by person and project, not by app."),
    (PEACH, "“Why is my laptop hot?”",
     "Names the process, shows what it cost you in CPU, RAM and battery, explains "
     "what it belongs to, and offers to throttle or kill it."),
    (GREEN, "“Reply to Arjun about Friday”",
     "Finds the thread in Telegram, reads the context, drafts in your voice, shows "
     "you the draft. Nothing sends until you approve it."),
    (YELLOW, "“Set me up for the demo”",
     "Launches the editor, dev server and browser, mutes non-urgent connectors, "
     "and restores the workspace afterwards."),
    (MAUVE, "It speaks first, when it matters",
     "A build failed, a deadline moved, a disk is nearly full. Proactive only inside "
     "the limits you set."),
], h=Inches(0.86))

# ======================================================= 05  architecture ====
s, y = slide_base("Five layers, one loop: perceive → reason → act",
                  "Architecture")
layers = [
    (MAUVE, "Interaction", "Overlay · global hotkey · voice · CLI · Plasma widget"),
    (BLUE, "Synapse Core", "Planner · memory · policy engine · model router"),
    (TEAL, "Perception", "Message bus · window & focus state · process table · file index"),
    (GREEN, "Action", "Connector API · process control · input synthesis · shell"),
    (OVERLAY, "Substrate", "Arch Linux · systemd · D-Bus · Wayland · Plasma 6"),
]
lh, lg = Inches(0.84), Inches(0.15)
for i, (cl, name, detail) in enumerate(layers):
    yy = y + i * (lh + lg)
    card(s, M, yy, CW - Inches(1.9), lh, MANTLE, cl)
    write(tbox(s, M + Inches(0.32), yy, Inches(2.5), lh, MSO_ANCHOR.MIDDLE),
          [{"t": name, "size": 15, "c": cl, "bold": True}], first=True)
    write(tbox(s, M + Inches(2.95), yy, CW - Inches(5.1), lh, MSO_ANCHOR.MIDDLE),
          [{"t": detail, "size": 11.5, "c": SUBTEXT}], first=True)

# side rail: the loop
rx = M + CW - Inches(1.72)
card(s, rx, y, Inches(1.72), 5 * lh + 4 * lg, MANTLE, None)
write(tbox(s, rx + Inches(0.2), y + Inches(0.26), Inches(1.32), Inches(4.2)),
      [{"t": "THE LOOP", "size": 9.5, "c": OVERLAY, "bold": True, "space_after": 12},
       {"t": "sense", "size": 12.5, "c": TEAL, "bold": True, "space_after": 3},
       {"t": "normalise", "size": 12.5, "c": SAPPHIRE, "bold": True, "space_after": 3},
       {"t": "reason", "size": 12.5, "c": BLUE, "bold": True, "space_after": 3},
       {"t": "check policy", "size": 12.5, "c": YELLOW, "bold": True, "space_after": 3},
       {"t": "act", "size": 12.5, "c": GREEN, "bold": True, "space_after": 3},
       {"t": "audit", "size": 12.5, "c": MAUVE, "bold": True}], first=True)

# ================================================== 06  core / the daemon ====
s, y = slide_base("Synapse Core: the reasoning daemon", "Component deep-dive")
write(tbox(s, M, y, Inches(10.8), Inches(0.42)),
      [{"t": "A user-level systemd service. Never root — privilege comes from "
             "explicit capabilities, not from uid 0.",
        "size": 14, "c": SUBTEXT, "line": 1.35}], first=True)
grid(s, y + Inches(0.62), [
    (BLUE, "Planner",
     "Turns intent into a typed action plan. Each step is validated against the "
     "policy engine before anything runs."),
    (MAUVE, "Memory",
     "Local store of entities, threads and preferences. Domain-scoped, so work "
     "context never leaks into personal."),
    (YELLOW, "Policy engine",
     "The gate every action passes through. Deny-by-default, scoped grants, "
     "confirm-before-destructive."),
    (TEAL, "Model router",
     "Routine triage runs on-device (llama.cpp / Ollama). Hard reasoning escalates "
     "to a cloud model only on opt-in."),
    (GREEN, "Agent backends",
     "claude-code, codex and opencode already ship in the image and plug in as "
     "pluggable executors for code tasks."),
    (SAPPHIRE, "Event scheduler",
     "Watches triggers — a failed build, a moved deadline, a filling disk — and "
     "decides whether to speak up."),
], cols=3, h=Inches(1.78))

# ========================================================= 07  perception ====
s, y = slide_base("How it sees the machine", "Perception layer")
rows(s, y, [
    (TEAL, "Message ingest",
     "Connector plugins normalise every account into one event schema: sender, "
     "thread, entities, urgency, sensitivity label."),
    (BLUE, "App & window state",
     "Wayland and Plasma report focus, window titles, idle time and session state — "
     "so “this file” and “this window” resolve to something real."),
    (SAPPHIRE, "Process table",
     "Continuous /proc and cgroup sampling: CPU, RSS, disk and network per PID, "
     "attributed back to the app that spawned it."),
    (GREEN, "Filesystem index",
     "Incremental inotify-backed index of paths and recent edits. Content is "
     "indexed locally, never uploaded by default."),
    (YELLOW, "System signals",
     "Battery, thermals, disk pressure, unit failures and network state from systemd "
     "and D-Bus."),
], h=Inches(0.9))

# ===================================================== 08  identity/login ====
s, y = slide_base("Logging into every app, without holding your passwords",
                  "Connectors & identity")
card(s, M, y, CW, Inches(0.98), MANTLE, RED)
write(tbox(s, M + Inches(0.42), y, CW - Inches(0.9), Inches(0.98), MSO_ANCHOR.MIDDLE),
      [{"t": "Design rule: the model never sees a credential. It requests a scoped "
             "action; the identity broker signs and performs the call.",
        "size": 15, "c": WHITE, "bold": True, "line": 1.25}], first=True)
grid(s, y + Inches(1.24), [
    (GREEN, "Identity broker",
     "OAuth tokens and session keys live in the kernel keyring, encrypted at rest. "
     "Out-of-process from the model, always."),
    (BLUE, "Tier 1 — official APIs",
     "Slack, Gmail, Telegram, Discord, GitHub. Real OAuth scopes, real revocation, "
     "no terms-of-service grey area."),
    (YELLOW, "Tier 2 — session bridge",
     "Apps with no API get a user-authorised local bridge, clearly disclosed and "
     "individually switchable."),
    (TEAL, "Scopes, per connector",
     "Read-only on connect. Sending, deleting and posting each need a separate "
     "explicit grant."),
    (MAUVE, "Revoke in one place",
     "A single panel lists every connector, scope and token, with a one-click "
     "disconnect and local data wipe."),
    (PEACH, "Fail closed",
     "A broker that cannot verify a scope refuses the call. No silent fallback to "
     "broader permission."),
], cols=3, h=Inches(1.74))

# ==================================================== 09  process control ====
s, y = slide_base("Running and killing processes, safely", "Action layer")
grid(s, y, [
    (GREEN, "Verbs it exposes",
     "launch · stop · kill · renice · cgroup-throttle · suspend · restart-unit. "
     "Each one is a distinct capability, granted separately."),
    (RED, "Protected set",
     "PID 1, the session, display server, keyring and the core itself are "
     "structurally unkillable — enforced by the broker, not by prompt wording."),
    (YELLOW, "Confirm, then remember",
     "Destructive actions ask once with a plain-language summary. Approve a pattern "
     "and it runs unattended next time."),
], cols=3, h=Inches(1.82))
grid(s, y + Inches(2.08), [
    (BLUE, "Explain before acting",
     "“chrome_crashpad, 180% CPU for 40 minutes, ~18% of your battery today.” "
     "You see the reasoning, then decide."),
    (TEAL, "Reversible where possible",
     "Prefer throttle over kill and suspend over terminate. Killing something "
     "unsaved requires an explicit override."),
    (MAUVE, "Everything audited",
     "Who asked, what was run, which capability authorised it, what changed. "
     "Append-only, readable by you."),
], cols=3, h=Inches(1.82))

# =============================================== 10  message intelligence ====
s, y = slide_base("Every message, one timeline", "Message intelligence")
rows(s, y, [
    (SAPPHIRE, "Unified stream",
     "All connectors collapse into a single chronological view, deduplicated across "
     "apps when the same thread arrives twice."),
    (PEACH, "Triage, not summarising",
     "Sorted by who is blocked, what has a deadline and what is genuinely FYI — "
     "rather than by arrival time."),
    (BLUE, "Entity linking",
     "People, projects, files and commits are resolved across apps, so a Slack "
     "thread and a Telegram DM about the same thing connect."),
    (GREEN, "Draft-only by default",
     "Replies are written in your voice and shown to you. Auto-send stays off until "
     "you enable it per contact."),
    (RED, "Message content is untrusted",
     "Incoming text is data, never instruction. A message asking the assistant to "
     "run something cannot cause a tool call."),
], h=Inches(0.9))

# ====================================================== 11  permissioning ====
s, y = slide_base("The user holds the leash", "Trust model")
modes = [
    (TEAL, "OBSERVE", "Perceives and answers. Zero side effects.",
     "Default on first boot"),
    (YELLOW, "ASSIST", "Drafts, stages and proposes. You confirm each effect.",
     "Everyday mode"),
    (PEACH, "ACT", "Executes approved patterns unattended, inside scope.",
     "Opt-in, per capability"),
]
cw3 = int((CW - Inches(0.44)) / 3)
for i, (cl, name, body, note) in enumerate(modes):
    x = M + i * (cw3 + Inches(0.22))
    card(s, x, y, cw3, Inches(1.86), MANTLE, cl)
    write(tbox(s, x + Inches(0.32), y + Inches(0.26), cw3 - Inches(0.6), Inches(1.4)),
          [{"t": name, "size": 17, "c": cl, "bold": True, "space_after": 8},
           {"t": body, "size": 12, "c": TEXT, "line": 1.3, "space_after": 8},
           {"t": note, "size": 10.5, "c": OVERLAY, "line": 1.2}], first=True)

rows(s, y + Inches(2.16), [
    (MAUVE, "Capability tokens",
     "Every grant is a tuple of resource, verb and duration. “Read Telegram for "
     "7 days” is a different token from “send as me”."),
    (BLUE, "Consent in plain words",
     "The prompt states the concrete effect, not a permission category: “Send this "
     "message to Arjun on Telegram now.”"),
    (RED, "Kill switch",
     "One shortcut suspends all perception and action instantly, and the state is "
     "visible in the panel at all times."),
], h=Inches(0.82))

# ======================================================= 12  privacy/risk ====
s, y = slide_base("Local-first, and honest about the blast radius",
                  "Privacy & security")
card(s, M, y, CW, Inches(0.9), MANTLE, RED)
write(tbox(s, M + Inches(0.42), y, CW - Inches(0.9), Inches(0.9), MSO_ANCHOR.MIDDLE),
      [{"t": "Stated threat model: an assistant with this much reach is a "
             "high-value target. A compromise here is a total compromise — so the "
             "architecture assumes it will be attacked.",
        "size": 14, "c": WHITE, "bold": True, "line": 1.25}], first=True)
grid(s, y + Inches(1.16), [
    (GREEN, "Data stays home",
     "Perception data never leaves the device by default. Cloud calls are per-task, "
     "opt-in, and shown to you when they happen."),
    (TEAL, "On-device inference",
     "Triage, classification and routine drafting run locally, so the common path "
     "involves no network at all."),
    (BLUE, "Sensitivity labels",
     "OTPs, keys and banking content are tagged at ingest and excluded from model "
     "context and from any remote call."),
    (YELLOW, "Least privilege",
     "The core runs unprivileged; connectors are sandboxed per-app, so one bad "
     "connector cannot reach another's data."),
    (MAUVE, "Readable audit log",
     "Append-only record of every perception source and every action, exportable "
     "and inspectable without special tooling."),
    (PEACH, "Injection resistance",
     "Untrusted content is fenced from the instruction channel; destructive verbs "
     "always require a human-confirmed capability."),
], cols=3, h=Inches(1.72))

# ============================================================== 13  stack ====
s, y = slide_base("Built on things that already work", "Technology")
write(tbox(s, M, y, Inches(10.8), Inches(0.42)),
      [{"t": "The base distribution exists today: an archiso profile shipping Plasma 6 "
             "and three agent CLIs. The AI layer is the next build on top.",
        "size": 14, "c": SUBTEXT, "line": 1.35}], first=True)
stack = [
    (OVERLAY, "Base", "Arch Linux · archiso · Calamares installer"),
    (BLUE, "Desktop", "Plasma 6 on Wayland · Catppuccin Macchiato · SDDM"),
    (TEAL, "Plumbing", "systemd user units · D-Bus IPC · cgroups v2 · inotify"),
    (MAUVE, "Core", "Rust daemon · Python connector SDK · SQLite + vector index"),
    (GREEN, "Models", "Ollama / llama.cpp local · cloud escalation on opt-in"),
    (YELLOW, "Agents", "claude-code · codex · opencode as pluggable executors"),
    (SAPPHIRE, "Secrets", "kernel keyring · per-connector sandbox · OAuth broker"),
    (PEACH, "Packaging", "pacman repo · reproducible ISO · signed connector plugins"),
]
cw2 = int((CW - Inches(0.22)) / 2)
for i, (cl, head, body) in enumerate(stack):
    x = M + (i % 2) * (cw2 + Inches(0.22))
    yy = y + Inches(0.66) + (i // 2) * Inches(1.02)
    card(s, x, yy, cw2, Inches(0.88), MANTLE, cl)
    write(tbox(s, x + Inches(0.3), yy, Inches(1.5), Inches(0.88), MSO_ANCHOR.MIDDLE),
          [{"t": head, "size": 13, "c": cl, "bold": True}], first=True)
    write(tbox(s, x + Inches(1.85), yy, cw2 - Inches(2.15), Inches(0.88),
               MSO_ANCHOR.MIDDLE),
          [{"t": body, "size": 11, "c": SUBTEXT, "line": 1.25}], first=True)

# ============================================================ 14  roadmap ====
s, y = slide_base("From distro to AI-native OS", "Roadmap")
phases = [
    (GREEN, "P0", "Foundation", "Shipping",
     "Arch + Plasma 6 image, Calamares installer, agent CLIs preinstalled."),
    (TEAL, "P1", "Perception", "Next",
     "Process and window telemetry, event schema, first two connectors, unified timeline."),
    (BLUE, "P2", "Action", "",
     "Capability engine, consent UI, process control, draft replies, audit log."),
    (MAUVE, "P3", "Autonomy", "",
     "Proactive briefings, multi-step plans, local model routing, learned approvals."),
    (YELLOW, "P4", "Ecosystem", "",
     "Public connector SDK, signed third-party plugins, org policy profiles."),
]
ph, pg = Inches(0.94), Inches(0.14)
for i, (cl, tag, name, badge, body) in enumerate(phases):
    yy = y + i * (ph + pg)
    card(s, M, yy, CW, ph, MANTLE, cl)
    rect(s, M + Inches(0.3), yy + Inches(0.235), Inches(0.47), Inches(0.47),
         SURFACE0, MSO_SHAPE.ROUNDED_RECTANGLE, 0.22)
    write(tbox(s, M + Inches(0.3), yy + Inches(0.355), Inches(0.47), Inches(0.24)),
          [{"t": tag, "size": 11.5, "c": cl, "bold": True, "font": MONO,
            "align": PP_ALIGN.CENTER}], first=True)
    write(tbox(s, M + Inches(0.95), yy, Inches(1.9), ph, MSO_ANCHOR.MIDDLE),
          [{"t": name, "size": 14, "c": WHITE, "bold": True}], first=True)
    write(tbox(s, M + Inches(2.95), yy, CW - Inches(4.6), ph, MSO_ANCHOR.MIDDLE),
          [{"t": body, "size": 11.5, "c": SUBTEXT, "line": 1.28}], first=True)
    if badge:
        bw = Inches(0.98)
        rect(s, M + CW - bw - Inches(0.3), yy + Inches(0.3), bw, Inches(0.32),
             SURFACE0, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5)
        write(tbox(s, M + CW - bw - Inches(0.3), yy + Inches(0.375), bw, Inches(0.22)),
              [{"t": badge, "size": 9.5, "c": cl, "bold": True,
                "align": PP_ALIGN.CENTER}], first=True)

# ============================================================== 15  risks ====
s, y = slide_base("What could go wrong, and the answer", "Risks & mitigations")
hh = Inches(0.42)
card(s, M, y, CW, hh, SURFACE0, None, 0.12)
write(tbox(s, M + Inches(0.32), y, Inches(4.4), hh, MSO_ANCHOR.MIDDLE),
      [{"t": "RISK", "size": 10, "c": SUBTEXT, "bold": True}], first=True)
write(tbox(s, M + Inches(5.1), y, Inches(7.0), hh, MSO_ANCHOR.MIDDLE),
      [{"t": "MITIGATION", "size": 10, "c": SUBTEXT, "bold": True}], first=True)
risks = [
    (RED, "Credential blast radius",
     "Broker-held tokens in the kernel keyring; the model never holds one; per-scope "
     "grants and one-click revocation."),
    (PEACH, "Prompt injection via messages",
     "Incoming content is data, not instruction. No tool call originates from message "
     "text; destructive verbs need human confirmation."),
    (YELLOW, "Hallucinated destructive action",
     "Deny-by-default capabilities, a structurally protected process set, reversible "
     "verbs preferred, undo where possible."),
    (BLUE, "Unofficial app bridges",
     "Official APIs first; bridges are opt-in, disclosed and individually "
     "switchable, with the trade-off stated plainly."),
    (TEAL, "Local model quality",
     "Route by difficulty: on-device for triage, escalate hard reasoning to a "
     "stronger model with explicit consent."),
    (MAUVE, "Always-on telemetry cost",
     "Adaptive sampling and idle backoff; the perception budget is visible and "
     "capped by the user."),
]
rh, rg = Inches(0.66), Inches(0.1)
for i, (cl, risk, fix) in enumerate(risks):
    yy = y + hh + Inches(0.12) + i * (rh + rg)
    card(s, M, yy, CW, rh, MANTLE, cl)
    write(tbox(s, M + Inches(0.32), yy, Inches(4.35), rh, MSO_ANCHOR.MIDDLE),
          [{"t": risk, "size": 12, "c": WHITE, "bold": True, "line": 1.15}], first=True)
    rect(s, M + Inches(4.85), yy + Inches(0.13), Inches(0.012), rh - Inches(0.26),
         SURFACE1)
    write(tbox(s, M + Inches(5.1), yy, CW - Inches(5.5), rh, MSO_ANCHOR.MIDDLE),
          [{"t": fix, "size": 11, "c": SUBTEXT, "line": 1.24}], first=True)

# ============================================================= 16  closing ====
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, CRUST)
rect(s, 0, 0, Inches(0.09), H, MAUVE)
write(tbox(s, M, Inches(1.5), Inches(11.2), Inches(0.32)),
      [{"t": "WHY NOW", "size": 12, "c": MAUVE, "bold": True}], first=True)
write(tbox(s, M, Inches(2.0), Inches(11.4), Inches(1.5)),
      [{"t": "The assistant should live in\nthe OS, not in a tab.",
        "size": 44, "c": WHITE, "bold": True, "line": 1.16}], first=True)
grid_items = [
    (TEAL, "Local models got good enough",
     "Triage and drafting now run on consumer hardware, which makes an always-on "
     "perception layer viable without shipping your life to a server."),
    (BLUE, "Capable agents already exist",
     "Tool-using coding agents ship in this image today. What they lack is a "
     "system-level view and a permission model."),
    (GREEN, "Arch gives full control",
     "No vendor assistant to work around. The distro owns the desktop, the services "
     "and the packaging end to end."),
]
cw3 = int((CW - Inches(0.44)) / 3)
for i, (cl, head, body) in enumerate(grid_items):
    x = M + i * (cw3 + Inches(0.22))
    card(s, x, Inches(4.0), cw3, Inches(1.92), MANTLE, cl)
    write(tbox(s, x + Inches(0.32), Inches(4.26), cw3 - Inches(0.62), Inches(1.5)),
          [{"t": head, "size": 14, "c": WHITE, "bold": True, "space_after": 8,
            "line": 1.1},
           {"t": body, "size": 11, "c": SUBTEXT, "line": 1.32}], first=True)
write(tbox(s, M, Inches(6.4), CW, Inches(0.3)),
      [{"t": "SynapseOS  ·  an Arch-based, AI-native operating system",
        "size": 11.5, "c": OVERLAY}], first=True)

out = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "SynapseOS-AI-Native-OS.pptx"
prs.save(str(out))
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
